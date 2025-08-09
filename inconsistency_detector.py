import pptx
import google.generativeai as genai
from PIL import Image
import os
from dotenv import load_dotenv
import glob
import argparse

# Load the API key from the .env file
load_dotenv()
api_key = os.getenv("GEMINI_API_KEY")
genai.configure(api_key=api_key)

def get_ppt_content(deck_path):
    
    all_slide_info = {}
    presentation = pptx.Presentation(deck_path)
    
    for i, slide in enumerate(presentation.slides):
        slide_text_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text_content.append(shape.text)
            elif shape.has_table:
                table = shape.table
                rows_as_text = []
                for row in table.rows:
                    row_cells = [cell.text for cell in row.cells]
                    rows_as_text.append(', '.join(row_cells))
                slide_text_content.append('Table content: ' + '; '.join(rows_as_text))
        
        all_slide_info[f"Slide {i+1}"] = ' '.join(slide_text_content).strip()
    return all_slide_info

def describe_images_with_gemini(image_paths):
    image_summaries={}
    model = genai.GenerativeModel('gemini-2.5-flash') 
    
    for path in image_paths:
        slide_image = Image.open(path)
        prompt = "Describe all text, charts, and key data points on this slide image."
        response = model.generate_content([prompt, slide_image])
        
        slide_name = os.path.basename(path).split('.')[0]
        image_summaries[slide_name] = response.text
        
    return image_summaries

def find_inconsistencies(all_slide_content):
    """Sends all presentation content to Gemini for a full inconsistency analysis."""
    model = genai.GenerativeModel('gemini-2.5-flash') # Updated here
    
    combined_content = "Here is the content from a multi-slide PowerPoint presentation:\n\n"
    for slide_name, content in all_slide_content.items():
        combined_content += f"--- {slide_name} ---\n{content}\n\n"
    
    prompt = (
        f"Analyze the following presentation for any factual or logical inconsistencies. "
        f"Look for conflicting numbers, contradictory claims, and timeline mismatches. "
        f"Provide a clear, structured output, referencing slide numbers and explaining the nature of each issue. "
        f"If no inconsistencies are found, state that clearly.\n\n"
        f"Presentation Content:\n{combined_content}"
    )

    response = model.generate_content(prompt)
    return response.text

def main():
    
    parser = argparse.ArgumentParser(description="Finds inconsistencies in a PowerPoint presentation.")
    parser.add_argument('pptx_path', type=str, help='Path to the PowerPoint (.pptx) file.')
    parser.add_argument('image_dir', type=str, help='Directory containing the slide images.')
    args = parser.parse_args()

    # get file paths from command-line arguments
    deck_path = args.pptx_path
    image_directory = args.image_dir

    # find all JPEG images in the specified directory
    image_paths = sorted(glob.glob(os.path.join(image_directory, '*.jpeg')))

    # Extract content from both the deck file and the images
    ppt_content = get_ppt_content(deck_path)
    image_descriptions = describe_images_with_gemini(image_paths)

    # Combine all content for a full analysis
    all_content_for_analysis = {**ppt_content, **image_descriptions}

    # run the analysis and get the report
    inconsistency_report = find_inconsistencies(all_content_for_analysis)

    print(inconsistency_report)

if __name__ == '__main__':
    main()